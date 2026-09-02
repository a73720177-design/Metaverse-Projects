from pathlib import Path

from pptx import Presentation


def parse_ppt(path: Path) -> list[tuple[int, str]]:
    presentation = Presentation(path)
    slides: list[tuple[int, str]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        slides.append((slide_number, "\n".join(parts)))
    return slides
