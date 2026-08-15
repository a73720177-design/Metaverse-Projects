import io
import os
from pptx import Presentation
import psycopg2
from minio import Minio
from dotenv import load_dotenv

# 1. 환경 변수 로드
load_dotenv()

# 2. MinIO 클라이언트 연결 (아까 구축한 정보)
minio_client = Minio(
    "25.15.210.79:9001",  # 또는 localhost
    access_key="minioadmin",
    secret_key="minioadmin",
    secure=False
)
bucket_name = "qwen-docs"

# 3. PostgreSQL 연결
db_conn = psycopg2.connect(
    host="25.15.210.79",  # 또는 localhost
    database="qwendb",
    user="admin",
    password="admin"
)
cursor = db_conn.cursor()

def process_pptx_from_minio(object_name, file_name):
    try:
        # A. MinIO에서 pptx 파일 다운로드 (메모리 상으로 가져오기)
        response = minio_client.get_object(bucket_name, object_name)
        pptx_data = io.BytesIO(response.read())
        
        # B. python-pptx로 파싱 시작
        prs = Presentation(pptx_data)
        
        extracted_texts = []
        for slide_index, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text += paragraph.text + " "
            
            if slide_text.strip():
                extracted_texts.append((slide_index + 1, slide_text.strip()))

        # C. PostgreSQL의 document_files에 파일 정보 저장
        minio_url = f"minio://{bucket_name}/{object_name}"
        cursor.execute(
            """
            INSERT INTO document_files (file_name, file_type, minio_url)
            VALUES (%s, %s, %s) RETURNING id;
            """,
            (file_name, "pptx", minio_url)
        )
        file_id = cursor.fetchone()[0]

        # D. 추출된 텍스트(슬라이드별 텍스트)를 document_chunks에 저장
        for slide_no, text in extracted_texts:
            chunk_content = f"[Slide {slide_no}] {text}"
            cursor.execute(
                """
                INSERT INTO document_chunks (file_id, chunk_text)
                VALUES (%s, %s);
                """,
                (file_id, chunk_content)
            )
        
        db_conn.commit()
        print(f"성공: {file_name} 파일이 파싱되어 DB에 저장되었습니다!")

    except Exception as e:
        db_conn.rollback()
        print(f"에러 발생: {e}")
    finally:
        cursor.close()
        db_conn.close()

# 실행 예시
# process_pptx_from_minio("sample.pptx", "발표자료.pptx")